# -*- coding: utf-8 -*-
"""
Generator presentasi PowerPoint (Bahasa Indonesia) untuk aplikasi
Bebang Sistem Informasi (BIS) — PT Prima Sarana Gemilang.

Menghasilkan: docs/presentasi/Presentasi-BIS.pptx
Menyisipkan screenshot dari docs/presentasi/screenshots/.

Jalankan:  python docs/presentasi/build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "screenshots")
OUT = os.path.join(HERE, "Presentasi-BIS.pptx")

# ---- Palet warna korporat (biru tua + aksen) ----
NAVY    = RGBColor(0x0F, 0x1E, 0x3D)   # latar gelap utama
BLUE    = RGBColor(0x1D, 0x4E, 0xD8)   # aksen biru (tombol app)
SKY     = RGBColor(0x3B, 0x82, 0xF6)   # biru terang
SLATE   = RGBColor(0x47, 0x55, 0x69)   # teks abu
LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)   # latar terang
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x0F, 0x17, 0x2A)   # teks gelap
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
CARD    = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------- Helper ----------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    """Latar penuh satu warna."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _set_fill(s, color)
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def band(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(s, color)
    s.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, text, size=18, color=INK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
        font="Segoe UI", line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=INK, gap=6,
            marker="—", marker_color=None, bold_lead=False):
    """Daftar butir. Item bisa str atau (judul, deskripsi)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    mc = marker_color or BLUE
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.1
        # marker
        rm = p.add_run(); rm.text = f"{marker}  "
        rm.font.size = Pt(size); rm.font.bold = True
        rm.font.color.rgb = mc; rm.font.name = "Segoe UI"
        if isinstance(it, tuple):
            head, desc = it
            rh = p.add_run(); rh.text = head + "  "
            rh.font.size = Pt(size); rh.font.bold = True
            rh.font.color.rgb = color; rh.font.name = "Segoe UI"
            rd = p.add_run(); rd.text = desc
            rd.font.size = Pt(size - 1); rd.font.color.rgb = SLATE
            rd.font.name = "Segoe UI"
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.bold = bold_lead
            r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb


def card(slide, x, y, w, h, fill=CARD, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.06
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def pic_framed(slide, img, x, y, w, h, border=SKY):
    """Bingkai + gambar (menjaga rasio 1440x900 = 1.6)."""
    ratio = 1440.0 / 900.0
    tw, th = w, Emu(int(w / ratio))
    if th > h:
        th = h; tw = Emu(int(h * ratio))
    ox = x + (w - tw) // 2
    oy = y + (h - th) // 2
    fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                ox - Emu(20000), oy - Emu(20000),
                                tw + Emu(40000), th + Emu(40000))
    _set_fill(fr, border); fr.shadow.inherit = False
    slide.shapes.add_picture(img, ox, oy, tw, th)


def footer(slide, page, dark=False):
    c = RGBColor(0x9C, 0xA3, 0xAF) if not dark else RGBColor(0x64, 0x74, 0x8B)
    txt(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
        "Bebang Sistem Informasi  •  PT Prima Sarana Gemilang", 10, c)
    txt(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.35),
        str(page), 10, c, align=PP_ALIGN.RIGHT)


def chip(slide, x, y, text, fill, tcolor=WHITE, w=Inches(1.9), h=Inches(0.42)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = tcolor; r.font.name = "Segoe UI"
    return s


# ---------------- Template slide ----------------
def content_header(slide, kicker, title):
    """Header konten: pita atas + kicker + judul."""
    band(slide, 0, 0, EMU_W, Inches(0.14), BLUE)
    txt(slide, Inches(0.6), Inches(0.35), Inches(11), Inches(0.35),
        kicker.upper(), 12, SKY, bold=True)
    txt(slide, Inches(0.6), Inches(0.68), Inches(12.1), Inches(0.8),
        title, 30, INK, bold=True)
    band(slide, Inches(0.6), Inches(1.5), Inches(1.1), Inches(0.06), BLUE)


def slide_shot(kicker, title, img, notes_items, caption=None, page=0):
    """Slide dua kolom: teks kiri, screenshot kanan (dibingkai)."""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    content_header(s, kicker, title)
    # kolom kiri: butir
    bullets(s, Inches(0.6), Inches(1.8), Inches(4.5), Inches(4.8),
            notes_items, size=15, gap=10)
    # kolom kanan: gambar
    pic_framed(s, img, Inches(5.3), Inches(1.75), Inches(7.5), Inches(4.7))
    if caption:
        txt(s, Inches(5.3), Inches(6.55), Inches(7.5), Inches(0.4),
            caption, 11, SLATE, italic=True, align=PP_ALIGN.CENTER)
    footer(s, page)
    return s


def P(name):
    return os.path.join(SHOTS, name)


# ================= SLIDE 1: SAMPUL =================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
# aksen dekoratif
band(s, 0, Inches(6.9), EMU_W, Inches(0.6), BLUE)
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(-1.2),
                       Inches(4.5), Inches(4.5))
_set_fill(d, RGBColor(0x14, 0x2A, 0x54)); d.shadow.inherit = False
d2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(4.2),
                        Inches(3.2), Inches(3.2))
_set_fill(d2, RGBColor(0x14, 0x2A, 0x54)); d2.shadow.inherit = False
# logo kotak
lg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5),
                        Inches(0.95), Inches(0.95))
lg.adjustments[0] = 0.2
_set_fill(lg, BLUE); lg.shadow.inherit = False
txt(s, Inches(0.9), Inches(1.6), Inches(0.95), Inches(0.8),
    "BIS", 22, WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, Inches(0.9), Inches(2.75), Inches(11), Inches(0.5),
    "PT PRIMA SARANA GEMILANG  •  IT DIVISION", 14, SKY, bold=True)
txt(s, Inches(0.85), Inches(3.2), Inches(11.5), Inches(1.4),
    "Bebang Sistem Informasi", 52, WHITE, bold=True)
txt(s, Inches(0.9), Inches(4.55), Inches(11), Inches(0.7),
    "Sistem Informasi Terintegrasi untuk HR, Inventory & Facility",
    22, RGBColor(0xCB, 0xD5, 0xE1))
txt(s, Inches(0.9), Inches(5.25), Inches(11), Inches(0.6),
    "Enterprise Resource Planning & Integrated Management System",
    15, RGBColor(0x94, 0xA3, 0xB8), italic=True)
chip(s, Inches(0.9), Inches(6.0), "HR", GREEN, w=Inches(1.3))
chip(s, Inches(2.35), Inches(6.0), "INVENTORY", BLUE, w=Inches(1.8))
chip(s, Inches(4.3), Inches(6.0), "FACILITY", SKY, w=Inches(1.7))
chip(s, Inches(6.15), Inches(6.0), "USER ACCESS", AMBER, w=Inches(2.0))
txt(s, Inches(0.9), Inches(7.0), Inches(9), Inches(0.4),
    "Presentasi Aplikasi  •  2026", 12, RGBColor(0xE2, 0xE8, 0xF0))


# ================= SLIDE 2: APA ITU BIS =================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
content_header(s, "Pengantar", "Apa itu Bebang Sistem Informasi?")
txt(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(1.1),
    "BIS adalah sistem informasi terintegrasi kelas enterprise yang "
    "menyatukan pengelolaan Sumber Daya Manusia, Inventaris/Logistik, dan "
    "Fasilitas dalam satu platform tunggal — dirancang untuk kebutuhan "
    "perusahaan pertambangan dan industri.", 17, SLATE, line_spacing=1.25)

cards = [
    ("groups", "Satu Sumber Data", "Data karyawan, aset, dan fasilitas "
     "terhubung — tidak ada lagi pencatatan ganda di banyak file terpisah.", GREEN),
    ("sync", "Proses Terstandar", "Alur persetujuan, audit, dan pelaporan "
     "mengikuti prosedur yang konsisten di seluruh unit kerja.", BLUE),
    ("shield", "Aman & Terkontrol", "Setiap akses dibatasi sesuai peran, "
     "seluruh perubahan tercatat dalam jejak audit.", AMBER),
]
cx = Inches(0.6); cw = Inches(3.9); gap = Inches(0.15)
for i, (ic, h, d, col) in enumerate(cards):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card(s, x, Inches(3.1), cw, Inches(3.0), fill=LIGHT)
    band(s, x, Inches(3.1), cw, Inches(0.12), col)
    txt(s, Emu(int(x)+Inches(0.3)), Inches(3.45), Inches(3.3), Inches(0.6),
        h, 18, INK, bold=True)
    txt(s, Emu(int(x)+Inches(0.3)), Inches(4.1), Inches(3.4), Inches(1.9),
        d, 14, SLATE, line_spacing=1.2)
footer(s, 2)


# ================= SLIDE 3: TANTANGAN & SOLUSI =================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
content_header(s, "Latar Belakang", "Tantangan yang Dijawab BIS")
# kolom kiri (masalah)
card(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(4.7), fill=RGBColor(0xFE,0xF2,0xF2))
band(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(0.12), RGBColor(0xDC,0x26,0x26))
txt(s, Inches(0.9), Inches(2.05), Inches(5.3), Inches(0.5),
    "Sebelum BIS", 20, RGBColor(0xB9,0x1C,0x1C), bold=True)
bullets(s, Inches(0.9), Inches(2.7), Inches(5.3), Inches(3.6), [
    "Data karyawan tersebar di banyak berkas Excel",
    "Aset perusahaan sulit dilacak — siapa memegang apa",
    "Stok gudang tidak real-time, rawan selisih",
    "Persetujuan manual, tidak ada jejak audit",
    "Laporan disusun manual, memakan waktu",
], size=15, gap=12, marker="✕", marker_color=RGBColor(0xDC,0x26,0x26), color=RGBColor(0x7F,0x1D,0x1D))
# kolom kanan (solusi)
card(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.7), fill=RGBColor(0xF0,0xFD,0xF4))
band(s, Inches(6.8), Inches(1.8), Inches(5.9), Inches(0.12), GREEN)
txt(s, Inches(7.1), Inches(2.05), Inches(5.3), Inches(0.5),
    "Dengan BIS", 20, RGBColor(0x15,0x80,0x3D), bold=True)
bullets(s, Inches(7.1), Inches(2.7), Inches(5.3), Inches(3.6), [
    "Basis data terpusat & terintegrasi",
    "Pelacakan aset per karyawan lewat serial/QR",
    "Stok & transaksi tercatat otomatis",
    "Alur persetujuan digital + jejak audit lengkap",
    "Laporan & ekspor Excel/PDF sekali klik",
], size=15, gap=12, marker="✓", marker_color=GREEN, color=RGBColor(0x14,0x53,0x2B))
footer(s, 3)


# ================= SLIDE 4: MODUL TERINTEGRASI =================
slide_shot(
    "Ikhtisar", "Empat Modul dalam Satu Platform",
    P("02-welcome-modul.png"),
    [
        ("Human Resources", "manajemen karyawan, master data, presensi & cuti."),
        ("Inventory", "produk, stok, transaksi, serial number & stock opname."),
        ("Facility Mgmt", "gedung, ruangan, akomodasi site & work order."),
        ("User Access", "keamanan, peran, dan hak akses pengguna."),
        ("Pemberitahuan sistem terpusat", "stok rendah, aset perlu ditinjau, persetujuan tertunda."),
    ],
    caption="Beranda Modul Utama — pintu masuk ke seluruh modul BIS",
    page=4,
)


# ================= SLIDE 5: ALUR KERJA / INTEGRASI =================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
content_header(s, "Cara Kerja", "Bagaimana Modul Saling Terhubung")
txt(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
    "Contoh: aset dari gudang diserahkan ke karyawan, ditempatkan di "
    "fasilitas, lalu diverifikasi lewat stock opname.", 15, SLATE)

steps = [
    ("inventory_2", "INVENTORY", "Barang masuk ke gudang dari supplier, "
     "dicatat dengan serial/tag.", BLUE),
    ("groups", "HR", "Aset diserahkan ke karyawan — tercatat siapa "
     "pemegangnya.", GREEN),
    ("hotel", "FACILITY", "Aset ditempatkan di gedung/ruangan site "
     "tertentu.", SKY),
    ("fact_check", "OPNAME & AUDIT", "Perhitungan fisik berkala + jejak "
     "audit seluruh perubahan.", AMBER),
]
bw = Inches(2.85); bh = Inches(2.9); by = Inches(2.7)
bx0 = Inches(0.6); gap = Inches(0.28)
for i, (ic, h, d, col) in enumerate(steps):
    x = Emu(int(bx0) + i * (int(bw) + int(gap)))
    card(s, x, by, bw, bh, fill=LIGHT)
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x)+Inches(0.3)),
                             Emu(int(by)+Inches(0.3)), Inches(0.7), Inches(0.7))
    _set_fill(num, col); num.shadow.inherit = False
    txt(s, Emu(int(x)+Inches(0.3)), Emu(int(by)+Inches(0.38)), Inches(0.7),
        Inches(0.55), str(i+1), 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Emu(int(x)+Inches(0.3)), Emu(int(by)+Inches(1.15)), Inches(2.3),
        Inches(0.5), h, 15, col, bold=True)
    txt(s, Emu(int(x)+Inches(0.3)), Emu(int(by)+Inches(1.6)), Inches(2.35),
        Inches(1.2), d, 12.5, SLATE, line_spacing=1.15)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON,
                                Emu(int(x)+int(bw)+Inches(0.02)),
                                Emu(int(by)+Inches(1.2)), Inches(0.24), Inches(0.5))
        _set_fill(ar, RGBColor(0xCB,0xD5,0xE1)); ar.shadow.inherit = False
card(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.75), fill=NAVY)
txt(s, Inches(0.9), Inches(6.02), Inches(11.6), Inches(0.55),
    "Satu aksi di satu modul otomatis memperbarui modul lain — tanpa input ganda.",
    15, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)


# ================= SLIDE 6: LOGIN & KEAMANAN AKSES =================
slide_shot(
    "Akses Masuk", "Login Aman Berbasis NIK",
    P("01-login.png"),
    [
        ("Autentikasi NIK + kata sandi", "kredensial terenkripsi (hash)."),
        ("Sesi berbasis token (JWT)", "aman & otomatis kedaluwarsa."),
        ("Opsi \"Ingat saya\"", "kemudahan pada perangkat tepercaya."),
        ("Antarmuka Bahasa Indonesia", "mudah dipahami seluruh pengguna."),
        ("Identitas perusahaan", "tampilan resmi PT Prima Sarana Gemilang."),
    ],
    caption="Halaman login BIS",
    page=6,
)


# ================= SLIDE 7: HR DASHBOARD EKSEKUTIF =================
slide_shot(
    "Modul HR", "Dashboard Eksekutif SDM",
    P("03-dashboard-hr.png"),
    [
        ("Ringkasan tenaga kerja", "total karyawan aktif dalam sekali pandang."),
        ("Sebaran per divisi/departemen", "komposisi organisasi tervisualisasi."),
        ("Grafik & statistik", "tren data SDM yang mudah dibaca pimpinan."),
        ("Data real-time", "angka selalu mengikuti kondisi terkini."),
        ("Dasar pengambilan keputusan", "informasi ringkas untuk manajemen."),
    ],
    caption="Dashboard HR — ikhtisar SDM untuk manajemen",
    page=7,
)


# ================= SLIDE 8: HR DIREKTORI KARYAWAN =================
slide_shot(
    "Modul HR", "Direktori & Data Karyawan",
    P("04-hr-karyawan.png"),
    [
        ("Data karyawan lengkap", "identitas, kepegawaian, keluarga & dokumen."),
        ("Pencarian & filter", "temukan karyawan secara cepat."),
        ("Foto & QR Code", "identitas visual tiap karyawan."),
        ("Impor & ekspor massal", "input banyak data via Excel."),
        ("Riwayat & jejak audit", "setiap perubahan tercatat rapi."),
    ],
    caption="Direktori Karyawan — pusat data seluruh pegawai",
    page=8,
)


# ================= SLIDE 9: HR MASTER DATA =================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
content_header(s, "Modul HR", "Master Data Kepegawaian Terstandar")
txt(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.8),
    "Seluruh pilihan data mengacu pada daftar baku yang seragam — mencegah "
    "penulisan tidak konsisten dan menjaga kualitas data organisasi.",
    16, SLATE, line_spacing=1.2)
md = [
    "Divisi", "Departemen", "Posisi / Jabatan", "Kategori Pangkat",
    "Golongan", "Sub Golongan", "Jenis Hubungan Kerja", "Tag",
    "Lokasi Kerja", "Status Karyawan",
]
gx = Inches(0.6); gy = Inches(2.9); cw = Inches(3.0); ch = Inches(0.95)
gxg = Inches(0.13); gyg = Inches(0.18)
for i, name in enumerate(md):
    r = i // 4; c = i % 4
    x = Emu(int(gx) + c * (int(cw) + int(gxg)))
    y = Emu(int(gy) + r * (int(ch) + int(gyg)))
    card(s, x, y, cw, ch, fill=LIGHT)
    band(s, x, y, Inches(0.1), ch, BLUE)
    txt(s, Emu(int(x)+Inches(0.35)), y, Emu(int(cw)-Inches(0.4)), ch,
        name, 15, INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
chip(s, Inches(0.6), Inches(6.35), "10 Entitas Master Data", NAVY, w=Inches(3.0))
chip(s, Inches(3.75), Inches(6.35), "Soft-delete & Kode Unik", SLATE, w=Inches(3.2))
footer(s, 9)


# ================= SLIDE 10: INVENTORY DASHBOARD =================
slide_shot(
    "Modul Inventory", "Dashboard Inventaris & Logistik",
    P("05-inventory-dashboard.png"),
    [
        ("Nilai & jumlah stok", "gambaran aset gudang secara ringkas."),
        ("Peringatan stok rendah", "notifikasi otomatis saat perlu pengadaan."),
        ("Aktivitas transaksi terkini", "pantau keluar-masuk barang."),
        ("Sebaran per gudang", "posisi stok di tiap lokasi."),
        ("Indikator kesehatan stok", "cepat kenali item bermasalah."),
    ],
    caption="Dashboard Inventory — kendali penuh atas logistik",
    page=10,
)


# ================= SLIDE 11: INVENTORY PRODUK & ASET =================
slide_shot(
    "Modul Inventory", "Katalog Produk & Aset",
    P("06-inventory-produk.png"),
    [
        ("Master produk lengkap", "kategori, sub-kategori, merek & satuan."),
        ("Serial number & Tag", "pelacakan unit aset secara individual."),
        ("Barang habis pakai (consumable)", "dukung stok pakai-habis."),
        ("Stok minimum", "ambang batas pemicu peringatan pengadaan."),
        ("Terhubung ke gudang", "posisi tiap produk selalu jelas."),
    ],
    caption="Katalog Produk — fondasi seluruh data inventaris",
    page=11,
)


# ================= SLIDE 12: INVENTORY STOK & TRANSAKSI =================
slide_shot(
    "Modul Inventory", "Transaksi Stok dengan Persetujuan",
    P("08-inventory-transaksi.png"),
    [
        ("Barang masuk & keluar", "seluruh pergerakan stok tercatat."),
        ("Alur persetujuan (approval)", "transaksi tervalidasi sebelum berlaku."),
        ("Void & amend", "koreksi terkontrol dengan jejak audit."),
        ("Retur & mutasi antar gudang", "dukung beragam skenario logistik."),
        ("Stok terupdate otomatis", "saldo selalu akurat setelah transaksi."),
    ],
    caption="Transaksi Inventory — pergerakan stok yang terkendali",
    page=12,
)


# ================= SLIDE 13: INVENTORY STOCK OPNAME =================
slide_shot(
    "Modul Inventory", "Stock Opname — Hitung Fisik Gudang",
    P("09-inventory-opname.png"),
    [
        ("Sesi hitung fisik terstruktur", "Draft → Berjalan → Selesai → Approved."),
        ("Snapshot stok & serial", "bandingkan fisik vs sistem otomatis."),
        ("Selisih dihitung otomatis", "temukan perbedaan tanpa hitung manual."),
        ("Penyesuaian saat disetujui", "koreksi stok resmi & tercatat."),
        ("Berita Acara PDF", "dokumen resmi hasil opname."),
    ],
    caption="Stock Opname — verifikasi fisik gudang berkala",
    page=13,
)


# ================= SLIDE 14: INVENTORY LABEL & QR =================
slide_shot(
    "Modul Inventory", "Label QR & Kartu Stok",
    P("10-inventory-label.png"),
    [
        ("Cetak label QR aset", "identitas fisik tiap unit barang."),
        ("Pindai cepat via QR", "cari data aset langsung dari lapangan."),
        ("Kartu stok per produk", "riwayat mutasi masuk-keluar lengkap."),
        ("Dukungan PWA/mobile", "pemindaian dari perangkat genggam."),
        ("Telusur aset akurat", "kurangi kehilangan & salah catat."),
    ],
    caption="Label QR — jembatan antara aset fisik & sistem",
    page=14,
)


# ================= SLIDE 15: FACILITY MANAGEMENT =================
slide_shot(
    "Modul Facility", "Manajemen Fasilitas & Akomodasi Site",
    P("13-facility-assets.png"),
    [
        ("Gedung & ruangan", "data lengkap fasilitas per lokasi."),
        ("Penghuni (occupants)", "siapa menempati ruang/mess site."),
        ("Aset fasilitas", "inventaris barang di tiap ruangan."),
        ("Work order pemeliharaan", "permintaan & pelacakan perbaikan."),
        ("Relevan untuk tambang", "kelola kamp & mess site terpencil."),
    ],
    caption="Facility Management — kelola fasilitas & akomodasi",
    page=15,
)


# ================= SLIDE 16: KEAMANAN & HAK AKSES (RBAC) =================
slide_shot(
    "Keamanan", "Kontrol Akses Berbasis Peran (RBAC)",
    P("15-admin-users.png"),
    [
        ("Manajemen user & akun", "kelola pengguna dari satu tempat."),
        ("Peran (role) fleksibel", "susun paket hak akses sesuai jabatan."),
        ("Hak akses granular", "13 sumber daya × 9 aksi terkontrol."),
        ("Akses per departemen", "batasi data sesuai lingkup kerja."),
        ("Jejak audit menyeluruh", "setiap tindakan penting terekam."),
    ],
    caption="Manajemen User & Peran — keamanan berlapis",
    page=16,
)


# ================= SLIDE 17: TEKNOLOGI & KEUNGGULAN =================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
content_header(s, "Di Balik Layar", "Teknologi & Keunggulan Sistem")
tech = [
    ("Modern & Andal", "Dibangun dengan teknologi web terkini (React, "
     "Node.js, PostgreSQL) — cepat, stabil, dan mudah dikembangkan."),
    ("Aman", "Enkripsi kata sandi, sesi token, kontrol akses berbasis "
     "peran, proteksi CSRF, dan pembatasan laju permintaan."),
    ("Cepat", "Caching cerdas (Redis) membuat halaman tampil responsif "
     "walau data besar."),
    ("Akses Mobile (PWA)", "Bisa dipasang seperti aplikasi & dukungan mode "
     "offline terbatas — cocok untuk operasional lapangan."),
    ("Terdokumentasi", "API terdokumentasi (Swagger) & teruji otomatis "
     "untuk menjaga kualitas."),
    ("Bahasa Indonesia", "Seluruh antarmuka & pesan dalam Bahasa Indonesia."),
]
tx = Inches(0.6); ty = Inches(1.85); cw = Inches(5.95); ch = Inches(1.45)
gxg = Inches(0.2); gyg = Inches(0.18)
for i, (h, d) in enumerate(tech):
    r = i // 2; c = i % 2
    x = Emu(int(tx) + c * (int(cw) + int(gxg)))
    y = Emu(int(ty) + r * (int(ch) + int(gyg)))
    card(s, x, y, cw, ch, fill=LIGHT)
    band(s, x, y, Inches(0.1), ch, SKY)
    txt(s, Emu(int(x)+Inches(0.35)), Emu(int(y)+Inches(0.15)),
        Emu(int(cw)-Inches(0.5)), Inches(0.5), h, 15, INK, bold=True)
    txt(s, Emu(int(x)+Inches(0.35)), Emu(int(y)+Inches(0.6)),
        Emu(int(cw)-Inches(0.5)), Inches(0.8), d, 11.5, SLATE, line_spacing=1.1)
footer(s, 17)


# ================= SLIDE 18: PENUTUP =================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
band(s, 0, 0, EMU_W, Inches(0.14), BLUE)
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.5),
                       Inches(4.5), Inches(4.5))
_set_fill(d, RGBColor(0x14, 0x2A, 0x54)); d.shadow.inherit = False
lg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.85), Inches(1.5),
                        Inches(1.6), Inches(1.6))
lg.adjustments[0] = 0.2
_set_fill(lg, BLUE); lg.shadow.inherit = False
txt(s, Inches(5.85), Inches(1.85), Inches(1.6), Inches(1.0),
    "BIS", 36, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.9),
    "Terima Kasih", 46, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.7),
    "Bebang Sistem Informasi — Satu Platform untuk HR, Inventory & Facility",
    18, RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
chip(s, Inches(4.4), Inches(5.5), "HR", GREEN, w=Inches(1.1))
chip(s, Inches(5.65), Inches(5.5), "INVENTORY", BLUE, w=Inches(1.6))
chip(s, Inches(7.4), Inches(5.5), "FACILITY", SKY, w=Inches(1.5))
txt(s, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
    "PT Prima Sarana Gemilang  •  IT Division  •  2026",
    13, RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)


# ================= SIMPAN =================
prs.save(OUT)
print("OK  ->", OUT)
print("Total slide:", len(prs.slides._sldIdLst))
